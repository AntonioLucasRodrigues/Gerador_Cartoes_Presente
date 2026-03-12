import pandas as pd
from pptx import Presentation
from copy import deepcopy
from pptx.util import Pt
from pptx.dml.color import RGBColor
import os

CAMINHO_EXCEL = "Relacao_Vales_Presente_Brinkelandia.xlsx" 
CAMINHO_MODELO_PPTX = "Cartao_Presente.pptx"
CAMINHO_OUTPUT_PPTX = "Cartoes_Prontos.pptx"

PLACEHOLDER_TEXTO = "{{NOME_PRESENTEADO}}"
COLUNA_NOME = 'Nomes para Vales Presente' 

def duplicate_slide(pres, slide_model, layout):
    new_slide = pres.slides.add_slide(layout)
    for shape in slide_model.shapes:
        el = shape.element
        new_el = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide
    
def gerar_cartoes_powerpoint():
    try:
        df_nomes = pd.read_excel(CAMINHO_EXCEL)
        lista_nomes = df_nomes[COLUNA_NOME].dropna().tolist()
    except FileNotFoundError:
        print(f"ERRO: Arquivo Excel não encontrado em: {CAMINHO_EXCEL}")
        return
    except KeyError:
        print(f"ERRO: Coluna '{COLUNA_NOME}' não encontrada no Excel.")
        return

    if not lista_nomes:
        print("A lista de nomes está vazia.")
        return

    try:
        prs = Presentation(CAMINHO_MODELO_PPTX)
    except FileNotFoundError:
        print(f"ERRO: Arquivo PowerPoint não encontrado em: {CAMINHO_MODELO_PPTX}")
        return

    slide_modelo = prs.slides[0]
    slide_layout_modelo = slide_modelo.slide_layout
    
    print(f"Modelo PPTX carregado. Encontrados {len(lista_nomes)} nomes no Excel.")

    for i, nome_cartao in enumerate(lista_nomes):
        novo_slide = duplicate_slide(prs, slide_modelo, slide_layout_modelo)
        
        for shape in novo_slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    if PLACEHOLDER_TEXTO in paragraph.text:
                        paragraph.text = paragraph.text.replace(PLACEHOLDER_TEXTO, nome_cartao)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = 'Montserrat Arabic' 
                            run.font.size = Pt(14) 
                            run.font.color.rgb = RGBColor(0, 0, 0)
                            run.font.bold = False 
                        break 
        
        print(f"Slide {i+1}/{len(lista_nomes)} gerado para: {nome_cartao}")

    def remove_slide(presentation, index):
        sldIdLst = presentation.slides._sldIdLst
        slide_element = sldIdLst[index]
        sldIdLst.remove(slide_element)
    
    try:
        remove_slide(prs, 0)
        print("Slide modelo original removido com sucesso.")
    except Exception as e:
        print(f"ERRO ao tentar remover o slide modelo. Remova o primeiro slide manualmente. Detalhe: {e}")
    
    prs.save(CAMINHO_OUTPUT_PPTX)
    print("\nProcesso concluído! 🎉")
    print(f"O arquivo '{CAMINHO_OUTPUT_PPTX}' foi criado.")

if __name__ == "__main__":
    gerar_cartoes_powerpoint()